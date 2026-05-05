from pathlib import Path
from .file_utils import list_files


def _extract_txt(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(file_path: Path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(file_path))
        pages = []
        for page in reader.pages:
            try:
                pages.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(pages)
    except ImportError:
        return f"[PDF extraction unavailable — install pypdf] {file_path.name}"
    except Exception as e:
        return f"[PDF read error: {e}]"


def _extract_docx(file_path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        return f"[DOCX extraction unavailable — install python-docx] {file_path.name}"
    except Exception as e:
        return f"[DOCX read error: {e}]"


def _extract_csv(file_path: Path) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(file_path, dtype=str)
        return df.to_string(index=False)
    except ImportError:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        return f"[CSV read error: {e}]"


def _extract_xlsx(file_path: Path) -> str:
    try:
        import pandas as pd
        sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)
        parts = []
        for sheet_name, df in sheets.items():
            parts.append(f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}")
        return "\n\n".join(parts)
    except ImportError:
        return f"[XLSX extraction unavailable — install pandas + openpyxl] {file_path.name}"
    except Exception as e:
        return f"[XLSX read error: {e}]"


def _extract_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    except ImportError:
        return f"[PPTX extraction unavailable — install python-pptx] {file_path.name}"
    except Exception as e:
        return f"[PPTX read error: {e}]"


def _extract_rtf(file_path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        raw = file_path.read_text(encoding="utf-8", errors="ignore")
        return rtf_to_text(raw)
    except ImportError:
        return f"[RTF extraction unavailable — install striprtf] {file_path.name}"
    except Exception as e:
        return f"[RTF read error: {e}]"


EXTRACTORS = {
    ".txt": _extract_txt,
    ".md": _extract_txt,
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".csv": _extract_csv,
    ".xlsx": _extract_xlsx,
    ".rtf": _extract_rtf,
    ".pptx": _extract_pptx,
}


CHUNK_SIZE = 1200
CHUNK_OVERLAP = 150


def extract_text(file_path: Path) -> str:
    """Extract text from a file, dispatching by extension."""
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if extractor is None:
        return f"[Unsupported file type: {ext}]"
    return extractor(file_path)


def _load_pdf_pages(file_path: Path) -> list:
    """Return one doc entry per PDF page so TF-IDF can score pages independently."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(file_path))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({
                    "path": str(file_path),
                    "text": text,
                    "source": file_path.name,
                    "chunk": i,
                })
        return pages
    except Exception:
        return []


def _split_chunks(text: str, file_path: Path) -> list:
    """Split long non-PDF text into overlapping chunks."""
    chunks = []
    start = 0
    idx = 0
    while start < len(text):
        chunk_text = text[start: start + CHUNK_SIZE]
        if chunk_text.strip():
            chunks.append({
                "path": str(file_path),
                "text": chunk_text,
                "source": file_path.name,
                "chunk": idx,
            })
        start += CHUNK_SIZE - CHUNK_OVERLAP
        idx += 1
    return chunks


def load_documents(data_dir: Path):
    data_dir = Path(data_dir)
    files = list_files(data_dir)
    docs = []
    for f in files:
        try:
            if f.suffix.lower() == ".pdf":
                docs.extend(_load_pdf_pages(f))
            else:
                text = extract_text(f)
                if len(text) > CHUNK_SIZE:
                    docs.extend(_split_chunks(text, f))
                else:
                    docs.append({"path": str(f), "text": text, "source": f.name})
        except Exception:
            continue
    return docs


def build_index(data_dir: Path):
    """Build/refresh document index from data_dir."""
    docs = load_documents(data_dir)
    return docs
